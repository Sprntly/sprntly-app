"""File conversion for uploaded dataset sources.

Mirrors the offline `scripts/convert_dataset.py` flow but exposes it as a
library so the upload endpoint can convert in-process. Each public converter
returns markdown as a string; the caller decides where to write it.

Rich converters exist for .docx, .xlsx, .csv, .pdf, .pptx, .html, .rtf, .txt,
.md. Any other type is *accepted, never rejected*: `convert` falls back to a
best-effort handler that passes the file through as text when it decodes
cleanly (yaml, json, toml, logs, …) or stores a placeholder stub for binary
content (audio, images, legacy .doc/.xls/.ppt, …). The raw bytes are always
saved by the caller, so processing can be improved later without changing the
upload contract. A .zip is not a converter here — it's expanded by
`app.datasets.ingest_zip` / the document-source upload route, which feed each
member back through these converters.

THE STUB IS MACHINE-DETECTABLE. A placeholder carries UNPARSED_STUB_MARKER and
`is_unparsed_stub()` recognizes it, because "we could not read this file" must
never be mistaken for content: it used to flow into the knowledge-graph
extractor, spending an LLM call to mine signals out of the sentence "its
content is not included in analysis yet", and — on the corpus path — get
recorded as permanently ingested, so the file was never retried even after a
parser for its type shipped. Ingestion paths skip stubs and leave them
unrecorded so a later parser picks them up.
"""
from __future__ import annotations

import io
import re
from pathlib import Path


# Lazy imports keep the dev environment lean and let tests stub modules out.
def _docx():
    import docx  # python-docx
    return docx


def _openpyxl():
    import openpyxl
    return openpyxl


def _pypdf():
    import pypdf
    return pypdf


def _pptx():
    from pptx import Presentation  # python-pptx
    return Presentation


def docx_to_md(data: bytes) -> str:
    doc = _docx().Document(io.BytesIO(data))
    parts: list[str] = []
    for p in doc.paragraphs:
        style = p.style.name if p.style else ""
        text = p.text
        if not text.strip():
            parts.append("")
            continue
        if style.startswith("Heading 1"):
            parts.append(f"# {text}")
        elif style.startswith("Heading 2"):
            parts.append(f"## {text}")
        elif style.startswith("Heading 3"):
            parts.append(f"### {text}")
        elif style.startswith("Heading"):
            parts.append(f"#### {text}")
        else:
            parts.append(text)
    for ti, table in enumerate(doc.tables):
        parts.append(f"\n[TABLE {ti}]")
        for row in table.rows:
            cells = [c.text.replace("\n", " ").strip() for c in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
    return "\n".join(parts)


#: Row ceilings for tabular formats. These are guards against a pathological
#: file, not a sampling policy — they used to be 30 (xlsx) and 200 (csv), which
#: silently discarded most of a real support/analytics/revenue export before
#: the model ever saw it, in exactly the evidence-bearing categories. The true
#: backstop downstream is document_sources.MAX_EXTRACTED_CHARS (200k), so these
#: can be generous.
XLSX_MAX_ROWS = 2000
CSV_MAX_ROWS = 5000


def xlsx_to_md(data: bytes, max_rows: int = XLSX_MAX_ROWS) -> str:
    wb = _openpyxl().load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"## {sheet}")
        rows = list(ws.iter_rows(values_only=True))
        truncated = False
        if len(rows) > max_rows:
            rows = rows[:max_rows]
            truncated = True
        for row in rows:
            cells = ["" if v is None else str(v).strip() for v in row]
            parts.append("| " + " | ".join(cells) + " |")
        if truncated:
            parts.append(f"\n_…(sheet truncated to {max_rows} rows)_")
        parts.append("")
    return "\n".join(parts)


def pdf_to_md(data: bytes) -> str:
    reader = _pypdf().PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"## Page {i + 1}\n\n{text.strip()}")
    return "\n\n".join(parts)


def pptx_to_md(data: bytes) -> str:
    """Extract slide text from a .pptx deck as markdown, one section per slide.

    Mirrors `pdf_to_md`: a `## Slide N` header per slide, then every shape's
    text (title, body bullets, text boxes) and any tables rendered as pipe
    rows. Diagram/image semantics are NOT recovered — decks are lossy, only the
    text layer is extracted (documented in the PRD-import contract). Empty
    slides are skipped so the output has no blank sections.
    """
    prs = _pptx()(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [c.text.replace("\n", " ").strip() for c in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                lines.append(text)
        if lines:
            parts.append(f"## Slide {i + 1}\n\n" + "\n".join(lines))
    return "\n\n".join(parts)


def txt_to_md(data: bytes) -> str:
    # plain text already qualifies as markdown for our LLM pipeline
    return data.decode("utf-8", errors="replace")


def html_to_md(data: bytes) -> str:
    """Readable text out of an HTML export, nav/script/style stripped.

    HTML decodes as text, so before this existed it "passed through" — meaning
    the extractor was handed raw tags and the page's navigation chrome as if it
    were customer evidence. Reuses the scraper's extractor so uploaded pages and
    scraped pages are read the same way.
    """
    from app.agents.scraper import extract_text

    return extract_text(data.decode("utf-8", errors="replace"))


# An RTF control word is a backslash, letters, an optional numeric parameter,
# and ONE optional trailing space that belongs to the control word, not the
# text. Matching that in a single pass matters: a two-pass approach that
# rewrites \par first would see `\parRenewal` (letters run on) and swallow the
# following word.
_RTF_CONTROL_RE = re.compile(r"\\([a-zA-Z]+)(-?\d+)?[ ]?")
_RTF_HEX_RE = re.compile(r"\\'[0-9a-fA-F]{2}")
_RTF_ESCAPED_CHAR_RE = re.compile(r"\\([{}\\])")

#: Control words that carry structure worth keeping as whitespace.
_RTF_BREAKS = {"par", "line", "sect", "page"}

#: Destination groups whose CONTENTS are metadata, not document text — font
#: and colour tables, styles, doc info, embedded pictures. Dropped wholesale;
#: everything else keeps its text.
_RTF_NOISE_DESTINATIONS = (
    "fonttbl", "colortbl", "stylesheet", "info", "pict", "generator",
    "themedata", "colorschememapping", "latentstyles", "datastore",
)


def _strip_rtf_groups(text: str) -> str:
    """Remove `{\\*...}` groups and known metadata destinations, brace-aware.

    A regex can't do this correctly — these groups nest (`{\\fonttbl{\\f0
    Times;}}`) — so this walks braces and skips a group whole once its opening
    control word marks it as noise.
    """
    out: list[str] = []
    depth = 0
    skip_until_depth: int | None = None
    i = 0
    while i < len(text):
        ch = text[i]
        if ch == "{":
            depth += 1
            if skip_until_depth is None:
                head = text[i + 1:i + 40]
                if head.startswith("\\*") or any(
                    head.startswith("\\" + d) for d in _RTF_NOISE_DESTINATIONS
                ):
                    skip_until_depth = depth
            i += 1
            continue
        if ch == "}":
            if skip_until_depth is not None and depth == skip_until_depth:
                skip_until_depth = None
            depth -= 1
            i += 1
            continue
        if skip_until_depth is None:
            out.append(ch)
        i += 1
    return "".join(out)


def rtf_to_md(data: bytes) -> str:
    """Plain text out of an RTF document.

    Like HTML, RTF decodes as text and so used to pass through raw — the model
    received `{\\rtf1\\ansi\\deff0…` control-word sludge as if it were customer
    evidence. This strips the control layer rather than adding a dependency:
    RTF is a dying format and a full parser isn't worth it. Formatting is
    discarded, words and paragraph breaks are kept.
    """
    text = data.decode("utf-8", errors="replace")
    text = _strip_rtf_groups(text)
    text = _RTF_HEX_RE.sub("", text)          # \'e9 escapes — drop, don't guess
    text = _RTF_CONTROL_RE.sub(
        lambda m: "\n" if m.group(1) in _RTF_BREAKS else "", text
    )
    text = _RTF_ESCAPED_CHAR_RE.sub(r"\1", text)  # literal \{ \} \\
    text = re.sub(r"\n{3,}", "\n\n", text)
    return "\n".join(line.strip() for line in text.splitlines()).strip()


def csv_to_md(data: bytes, max_rows: int = CSV_MAX_ROWS) -> str:
    """Convert CSV data to a markdown table."""
    import csv as _csv

    text = data.decode("utf-8", errors="replace")
    reader = _csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return "_Empty CSV file._\n"

    parts: list[str] = []
    header = rows[0]
    parts.append("| " + " | ".join(h.strip() for h in header) + " |")
    parts.append("| " + " | ".join("---" for _ in header) + " |")

    data_rows = rows[1:]
    truncated = len(data_rows) > max_rows
    for row in data_rows[:max_rows]:
        # Pad or trim row to match header length
        cells = row + [""] * (len(header) - len(row))
        parts.append("| " + " | ".join(c.strip() for c in cells[:len(header)]) + " |")

    if truncated:
        parts.append(f"\n_…(truncated to {max_rows} of {len(data_rows)} rows)_")

    return "\n".join(parts) + "\n"


# Routing -----------------------------------------------------------------

_SUFFIX_TO_CONVERTER = {
    ".docx": docx_to_md,
    ".xlsx": xlsx_to_md,
    ".csv": csv_to_md,
    # .tsv deliberately absent: csv.reader would split on commas and collapse
    # tab-separated data into one column. The textual fallback passes it
    # through intact, which reads better.
    ".pdf": pdf_to_md,
    ".pptx": pptx_to_md,
    ".html": html_to_md,
    ".htm": html_to_md,
    ".rtf": rtf_to_md,
    ".txt": txt_to_md,
    ".md": txt_to_md,  # already markdown, passthrough
}

SUPPORTED_SUFFIXES = tuple(_SUFFIX_TO_CONVERTER.keys())


class UnsupportedFileType(ValueError):
    """Retained for callers that still catch it; `convert` no longer raises it."""


def _looks_textual(data: bytes) -> bool:
    """Heuristic: does this look like decodable text rather than binary?"""
    if not data:
        return True
    if b"\x00" in data[:8192]:  # NUL byte → almost certainly binary
        return False
    try:
        data.decode("utf-8")
        return True
    except UnicodeDecodeError:
        return False


#: Marks a placeholder produced for a file we could not read. Kept as an HTML
#: comment so it is invisible wherever the markdown is rendered, and stable so
#: `is_unparsed_stub` can recognize stubs written by earlier releases too (see
#: the module docstring for why this must be detectable).
UNPARSED_STUB_MARKER = "<!-- sprntly:unparsed -->"

# Prose prefix inside the same placeholder, and the only marker stubs carried
# before the HTML comment shipped. `fallback_to_md` still writes it and
# `is_unparsed_stub` still matches it, so the stub text and its detector can
# never drift apart — and stubs already stored stay detectable.
_UNPARSED_STUB_MARKER = "_Stored as a source but not yet parsed (type "


def is_unparsed_stub(text: str | None) -> bool:
    """True iff `text` is a placeholder for a file we could not read.

    Callers use this to keep "we couldn't read this" out of anything that
    treats text as content — extraction, the corpus seed, evidence counts.

    The stub is deliberately NON-EMPTY so the file still registers as a source,
    which means an `if not text` check does NOT catch it. Anything that treats
    "no usable text" as a decision point (e.g. roadmap→KG ingest, which must not
    let an unparseable re-upload retire the previous roadmap's signals) has to
    ask this instead.
    """
    if not text:
        return False
    if UNPARSED_STUB_MARKER in text:
        return True
    # Legacy stubs (written before the marker existed).
    return _UNPARSED_STUB_MARKER in text


def fallback_to_md(filename: str, data: bytes) -> str:
    """Best-effort conversion for types without a dedicated converter.

    Never raises: textual content (yaml/json/toml/logs/…) passes through as-is;
    binary content (audio/images/legacy Office formats/…) becomes a placeholder
    stub so the file still lands as a source. The raw bytes are preserved by
    the caller either way, so a future parser can read the original.
    """
    if _looks_textual(data):
        return data.decode("utf-8", errors="replace")
    suffix = Path(filename).suffix.lower() or "(none)"
    kb = max(1, round(len(data) / 1024))
    return (
        f"{UNPARSED_STUB_MARKER}\n"
        f"# {Path(filename).name}\n\n"
        f"{_UNPARSED_STUB_MARKER}{suffix}, {kb} KB). "
        "Binary or unrecognized format — its content is not included in "
        "analysis yet._\n"
    )


def strip_nul(text: str) -> str:
    """Drop NUL characters from extracted text. THE ROOT-CAUSE FIX for a 500.

    `U+0000` IS VALID UTF-8, so a UTF-16LE file saved without a BOM, a PDF with
    padded strings, or a .docx carrying embedded binary all decode cleanly and
    pass every character-level check — and then Postgres refuses the write,
    because `text` cannot store a NUL. It answers SQLSTATE 22P05 ("unsupported
    Unicode escape sequence — \\u0000 cannot be converted to text") and PostgREST
    surfaces it as an APIError, which reaches the user as a 500 on a send.

    Observed live: a chat message with a document attached inlines the extracted
    text into the question, `db.asks.start_ask_job` inserts it, and the entire
    `POST /v1/ask` dies — not the attachment, the whole message.

    NO TEST IN THIS SUITE CAN CATCH THAT DOWNSTREAM. The SQLite fake stores NULs
    happily, so the failure only exists against real Postgres — the same trap
    `artifact_templates.store.TemplateSourceNotText` documents, which is why
    that path refuses a NUL outright at upload.

    Here it STRIPS rather than refuses, and the difference is deliberate: an
    uploaded FORMAT with a NUL is a file the user chose wrong and can re-save,
    while an attachment in mid-conversation is a document they want read. The
    NUL carries no meaning in either case — dropping it costs nothing and keeps
    the message working.
    """
    return text.replace("\x00", "") if text else text


def convert(filename: str, data: bytes) -> str:
    """Dispatch by extension; unknown types fall back instead of failing.

    Every converter's output passes through `strip_nul`, so no extraction path
    can hand a NUL to a caller that will try to persist it — see there."""
    suffix = Path(filename).suffix.lower()
    fn = _SUFFIX_TO_CONVERTER.get(suffix)
    if fn is None:
        return strip_nul(fallback_to_md(filename, data))
    return strip_nul(fn(data))


# Naming helpers ----------------------------------------------------------

_SLUG_RE = re.compile(r"[^a-z0-9_]+")


def slugify(name: str) -> str:
    """Normalize a filename stem into a safe markdown filename component.

    Lowercases, replaces non-alphanumerics with underscores, collapses runs,
    strips leading/trailing underscores. Empty input becomes 'untitled'.
    """
    s = _SLUG_RE.sub("_", name.strip().lower())
    s = re.sub(r"_+", "_", s).strip("_")
    return s or "untitled"


def md_filename(original_filename: str) -> str:
    """Pick the markdown filename for a converted upload. Preserves the stem."""
    stem = Path(original_filename).stem
    return f"{slugify(stem)}.md"
