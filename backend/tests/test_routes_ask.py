"""Tests for app.routes.ask — POST /v1/ask + GET /v1/ask/{ask_id}.

The Ask flow is fire-and-forget (blur/remount-safe), mirroring PRD/evidence:
POST persists a `generating` row in `ask_jobs` and kicks `qa_agent.answer(...)`
in a background task, returning `{ask_id, status}`; the client polls
GET /v1/ask/{ask_id} for the answer. A cache hit is persisted onto an
immediately-`ready` job so the POST contract is uniform; the user-visible
result is identical to the old synchronous body.

The route sits behind `require_company` and the `dataset` slug must resolve to
the caller's company — otherwise an arbitrary client slug would seed a FOREIGN
company's corpus into the LLM answer.

Key behaviours covered:
- auth gate (no session → 401)
- foreign / unowned dataset → 404 (corpus-leak denial)
- POST returns an ask_id + persists a generating job; the worker fills the
  answer; GET walks generating → ready with the SAME citation-stripped shape
- cache hit → ready job carrying the cached payload, no LLM call
- a worker failure marks the job `error` and never crashes
- GET for a foreign / nonexistent ask_id → 404
"""
from __future__ import annotations

import asyncio
import json
import time

from app import db
from app.routes import ask as ask_route


def _seed_corpus(data_dir, dataset, body="some corpus body"):
    ds = data_dir / dataset
    ds.mkdir(exist_ok=True)
    (ds / "a.md").write_text(body)


def _poll_ask(client, ask_id, *, timeout=5.0):
    """Poll GET /v1/ask/{id} until the job leaves `generating` (the worker runs
    on the TestClient's event loop, so the answer lands within a tick or two)."""
    deadline = time.monotonic() + timeout
    body = None
    while time.monotonic() < deadline:
        resp = client.get(f"/v1/ask/{ask_id}")
        assert resp.status_code == 200, resp.text
        body = resp.json()
        if body["status"] != "generating":
            return body
        time.sleep(0.02)
    return body


# ---- auth + tenant gate -----------------------------------------------------

def test_ask_without_session_returns_401(unauth_client, isolated_settings):
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    resp = unauth_client.post(
        "/v1/ask",
        json={"question": "What is the biggest churn driver?", "dataset": "acme"},
    )
    assert resp.status_code == 401


def test_ask_foreign_dataset_returns_404(tenant_client, isolated_settings):
    """Ask must reject a dataset slug that isn't the caller's — no corpus leak."""
    tenant_client.make(slug="company-a")
    _seed_corpus(isolated_settings["data_dir"], dataset="company-a")
    b = tenant_client.make(slug="company-b")
    resp = b.client.post(
        "/v1/ask",
        json={"question": "Tell me A's secrets?", "dataset": "company-a"},
    )
    assert resp.status_code == 404


def test_ask_unowned_dataset_returns_404(tenant_client, isolated_settings):
    """A slug that maps to no company at all is also 404 (fail closed)."""
    t = tenant_client.make(slug="acme")
    resp = t.client.post(
        "/v1/ask", json={"question": "What about nobody?", "dataset": "ghost"}
    )
    assert resp.status_code == 404


# ---- fire-and-forget contract (cache miss → worker) -------------------------

def test_ask_returns_ask_id_and_persists_generating_job(
    tenant_client, isolated_settings, fake_llm
):
    """POST returns {ask_id, status} and persists a job row for the caller."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    fake_llm["payload"] = {
        "answer": "## Finding\n\nThe answer.",
        "key_points": ["k1"],
        "citations": [{"source": "a", "evidence": "x"}],
        "confidence": 0.9,
        "unanswered": "",
    }
    resp = t.client.post(
        "/v1/ask",
        json={"question": "What is the biggest churn driver?", "dataset": "acme"},
    )
    assert resp.status_code == 200
    body = resp.json()
    assert isinstance(body["ask_id"], int)
    assert body["status"] in ("generating", "ready")
    # The job row exists and belongs to the caller's company.
    row = db.get_ask_job(body["ask_id"])
    assert row is not None
    assert row["company_id"] == t.company_id


def test_ask_worker_fills_answer_and_get_returns_old_shape(
    tenant_client, isolated_settings, fake_llm
):
    """The worker runs qa_agent.answer; GET walks generating → ready and returns
    the SAME citation-stripped shape the old synchronous POST returned."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    fake_llm["payload"] = {
        "answer": "## Finding\n\nThe answer.",
        "key_points": ["k1"],
        "citations": [{"source": "a", "evidence": "x"}],
        "confidence": 0.9,
        "unanswered": "",
    }
    start = t.client.post(
        "/v1/ask",
        json={"question": "What is the biggest churn driver?", "dataset": "acme"},
    ).json()
    body = _poll_ask(t.client, start["ask_id"])
    assert body["status"] == "ready"
    assert body["answer"] == "## Finding\n\nThe answer."
    # Citations stripped on the way out (same as the old endpoint).
    assert body["citations"] == []
    assert body["key_points"] == ["k1"]
    assert body["confidence"] == 0.9


def test_ask_cache_miss_invokes_fake_llm_exactly_once(
    tenant_client, isolated_settings, fake_llm
):
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    fake_llm["payload"] = {
        "answer": "x", "key_points": [], "citations": [],
        "confidence": 0.5, "unanswered": "",
    }
    start = t.client.post(
        "/v1/ask", json={"question": "Some unique question?", "dataset": "acme"}
    ).json()
    _poll_ask(t.client, start["ask_id"])
    assert len(fake_llm["calls"]) == 1


def test_ask_short_question_is_rejected(tenant_client, isolated_settings):
    """Pydantic min_length=3 — anything shorter is a validation error."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    resp = t.client.post("/v1/ask", json={"question": "hi", "dataset": "acme"})
    assert resp.status_code == 422


# ---- cache hit path ---------------------------------------------------------

def test_ask_cache_hit_returns_ready_job_without_llm_call(
    tenant_client, isolated_settings, fake_llm, monkeypatch
):
    monkeypatch.setattr(ask_route, "CACHE_HIT_DELAY_MIN_SECONDS", 0.0)
    monkeypatch.setattr(ask_route, "CACHE_HIT_DELAY_MAX_SECONDS", 0.0)

    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    question = "What are the biggest revenue drivers"
    cached_payload = {
        "answer": "**Cached answer**", "key_points": ["pre-warmed"],
        "citations": [], "confidence": 1.0, "unanswered": "",
    }
    cache_id = db.start_cached_ask(dataset="acme", question=question)
    db.complete_cached_ask(cache_id, json.dumps(cached_payload))

    fake_llm["calls"].clear()
    resp = t.client.post("/v1/ask", json={"question": question, "dataset": "acme"})
    assert resp.status_code == 200
    start = resp.json()
    # Cache hits resolve to a ready job immediately (no worker, no LLM).
    assert start["status"] == "ready"
    body = t.client.get(f"/v1/ask/{start['ask_id']}").json()
    assert body["status"] == "ready"
    assert body["answer"] == "**Cached answer**"
    assert body["key_points"] == ["pre-warmed"]
    assert fake_llm["calls"] == []


# ---- worker failure ---------------------------------------------------------

def test_ask_worker_failure_marks_error_and_does_not_crash(
    tenant_client, isolated_settings, fake_llm, monkeypatch
):
    """A failure inside the answer pipeline marks the job `error` (best-effort)
    and the worker never propagates the exception."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")

    import app.qa_agent as qa_agent_mod

    def _boom(**kwargs):  # noqa: ARG001
        raise RuntimeError("kaboom")

    monkeypatch.setattr(qa_agent_mod, "answer", _boom)

    start = t.client.post(
        "/v1/ask", json={"question": "A question that explodes?", "dataset": "acme"}
    ).json()
    body = _poll_ask(t.client, start["ask_id"])
    assert body["status"] == "error"
    assert "kaboom" in (body["error"] or "")


# ---- POST /v1/ask/{id}/cancel (stop an in-flight ask) -----------------------
# The composer's Stop button POSTs here. It flips a `generating` job to
# `cancelled`; the worker polls that status between LLM steps and aborts before
# the (expensive) answer call, and a late answer is discarded (complete/fail are
# guarded on status == 'generating'). Idempotent + race-safe + tenant-scoped.


def test_cancel_generating_job_flips_to_cancelled(tenant_client, isolated_settings):
    """A generating job → cancelled; GET reflects it."""
    t = tenant_client.make(slug="acme")
    ask_id = db.start_ask_job(company_id=t.company_id, dataset="acme", question="q?")
    resp = t.client.post(f"/v1/ask/{ask_id}/cancel")
    assert resp.status_code == 200, resp.text
    assert resp.json()["status"] == "cancelled"
    assert db.get_ask_job(ask_id)["status"] == "cancelled"
    assert t.client.get(f"/v1/ask/{ask_id}").json()["status"] == "cancelled"


def test_cancel_already_finished_job_is_noop_returns_terminal(
    tenant_client, isolated_settings
):
    """If the worker already finished (ready), cancel is a race-safe no-op and
    reports the real terminal status instead of clobbering it to cancelled."""
    t = tenant_client.make(slug="acme")
    ask_id = db.start_ask_job(company_id=t.company_id, dataset="acme", question="q?")
    db.complete_ask_job(ask_id, {
        "answer": "done", "key_points": [], "citations": [],
        "confidence": 1.0, "unanswered": "",
    })
    resp = t.client.post(f"/v1/ask/{ask_id}/cancel")
    assert resp.status_code == 200
    assert resp.json()["status"] == "ready"
    assert db.get_ask_job(ask_id)["status"] == "ready"


def test_cancel_foreign_job_returns_404_and_leaves_it_running(
    tenant_client, isolated_settings
):
    """A job belonging to another company is not cancellable (404, no
    disclosure) and stays generating."""
    a = tenant_client.make(slug="company-a")
    ask_id = db.start_ask_job(company_id=a.company_id, dataset="company-a", question="q?")
    b = tenant_client.make(slug="company-b")
    resp = b.client.post(f"/v1/ask/{ask_id}/cancel")
    assert resp.status_code == 404
    assert db.get_ask_job(ask_id)["status"] == "generating"


def test_cancel_nonexistent_returns_404(tenant_client, isolated_settings):
    t = tenant_client.make(slug="acme")
    resp = t.client.post("/v1/ask/999999/cancel")
    assert resp.status_code == 404


def test_late_answer_does_not_resurrect_a_cancelled_job(
    tenant_client, isolated_settings
):
    """The un-interruptible final LLM call can finish AFTER a cancel lands. The
    guarded complete_ask_job must NOT overwrite the cancel and resurface the
    unwanted answer."""
    t = tenant_client.make(slug="acme")
    ask_id = db.start_ask_job(company_id=t.company_id, dataset="acme", question="q?")
    assert db.cancel_ask_job(ask_id) == "cancelled"
    db.complete_ask_job(ask_id, {
        "answer": "late unwanted answer", "key_points": [], "citations": [],
        "confidence": 1.0, "unanswered": "",
    })
    row = db.get_ask_job(ask_id)
    assert row["status"] == "cancelled"
    assert (row.get("response") or {}).get("answer", "") != "late unwanted answer"


def test_worker_aborts_before_llm_when_cancelled(
    tenant_client, isolated_settings, fake_llm
):
    """A pre-cancelled job short-circuits at qa_agent's first checkpoint: the
    worker leaves the row `cancelled` (NOT error) and never calls the answer
    LLM, so the Stop actually saves the generation cost."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    ask_id = db.start_ask_job(company_id=t.company_id, dataset="acme", question="q?")
    db.cancel_ask_job(ask_id)
    fake_llm["calls"].clear()

    from app.ask_job_runner import run_ask_job

    asyncio.run(run_ask_job(
        ask_id=ask_id,
        enterprise_id=t.company_id,
        question="q?",
        dataset="acme",
    ))

    row = db.get_ask_job(ask_id)
    assert row["status"] == "cancelled"  # not 'error', not 'ready'
    assert fake_llm["calls"] == []       # the expensive answer LLM was skipped


# ---- GET status auth/ownership ----------------------------------------------

def test_get_ask_nonexistent_returns_404(tenant_client, isolated_settings):
    t = tenant_client.make(slug="acme")
    resp = t.client.get("/v1/ask/999999")
    assert resp.status_code == 404


def test_get_ask_foreign_job_returns_404(tenant_client, isolated_settings, fake_llm):
    """A job belonging to another company is not readable (404, no disclosure)."""
    a = tenant_client.make(slug="company-a")
    _seed_corpus(isolated_settings["data_dir"], dataset="company-a")
    fake_llm["payload"] = {
        "answer": "x", "key_points": [], "citations": [],
        "confidence": 0.5, "unanswered": "",
    }
    start = a.client.post(
        "/v1/ask", json={"question": "A's private question?", "dataset": "company-a"}
    ).json()
    b = tenant_client.make(slug="company-b")
    resp = b.client.get(f"/v1/ask/{start['ask_id']}")
    assert resp.status_code == 404


# ---- POST /v1/ask/extract-file ----------------------------------------------
# Parses a binary chat attachment (pptx/pdf/docx/…) to markdown so the composer
# can inline it as [Attached files] context — the fix for pptx attachments
# being silently dropped when sent with a plain question.

def _tiny_pptx(*, slides: bool = True) -> bytes:
    import io

    from pptx import Presentation

    prs = Presentation()
    if slides:
        s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        s.shapes.title.text = "Fraznet Enhancements"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_file_without_session_returns_401(unauth_client):
    resp = unauth_client.post(
        "/v1/ask/extract-file", files={"file": ("a.txt", b"hello", "text/plain")}
    )
    assert resp.status_code == 401


def test_extract_file_txt_passthrough(tenant_client):
    t = tenant_client.make(slug="acme")
    resp = t.client.post(
        "/v1/ask/extract-file",
        files={"file": ("notes.txt", b"churn is up 20%", "text/plain")},
    )
    assert resp.status_code == 200, resp.text
    body = resp.json()
    assert body["name"] == "notes.txt"
    assert "churn is up 20%" in body["markdown"]


def test_extract_file_pptx_returns_slide_markdown(tenant_client):
    import pytest

    pytest.importorskip("pptx")
    t = tenant_client.make(slug="acme")
    resp = t.client.post(
        "/v1/ask/extract-file",
        files={
            "file": (
                "deck.pptx",
                _tiny_pptx(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert resp.status_code == 200, resp.text
    body = resp.json()
    assert body["name"] == "deck.pptx"
    assert "## Slide 1" in body["markdown"]
    assert "Fraznet Enhancements" in body["markdown"]


def test_extract_file_empty_returns_400(tenant_client):
    t = tenant_client.make(slug="acme")
    resp = t.client.post(
        "/v1/ask/extract-file", files={"file": ("a.pptx", b"", "application/octet-stream")}
    )
    assert resp.status_code == 400


def test_extract_file_no_extractable_text_returns_422(tenant_client):
    import pytest

    pytest.importorskip("pptx")
    t = tenant_client.make(slug="acme")
    resp = t.client.post(
        "/v1/ask/extract-file",
        files={
            "file": (
                "empty-deck.pptx",
                _tiny_pptx(slides=False),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert resp.status_code == 422


def test_extract_file_over_size_cap_returns_413(tenant_client, monkeypatch):
    monkeypatch.setattr(ask_route, "_MAX_EXTRACT_BYTES", 10)
    t = tenant_client.make(slug="acme")
    resp = t.client.post(
        "/v1/ask/extract-file", files={"file": ("big.txt", b"x" * 11, "text/plain")}
    )
    assert resp.status_code == 413


def test_ask_accepts_question_with_inlined_attachment_block(tenant_client, isolated_settings, fake_llm):
    """The composer inlines extracted document markdown into the question
    ([Attached files] block) — the old 2000-char cap rejected any real deck."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    fake_llm["payload"] = {
        "answer": "ok", "key_points": [], "citations": [],
        "confidence": 0.9, "unanswered": "",
    }
    question = (
        "What are the riskiest requirements in this deck?\n\n[Attached files]\n"
        "--- deck.pptx ---\n" + ("Slide content line.\n" * 500)  # ~10k chars
    )
    resp = t.client.post("/v1/ask", json={"question": question, "dataset": "acme"})
    assert resp.status_code == 200, resp.text
    body = _poll_ask(t.client, resp.json()["ask_id"])
    assert body["status"] == "ready"


def test_find_cached_ask_skips_oversized_question_without_hitting_the_client():
    """An oversized question (a chat ask carrying an inlined [Attached files]
    block — tens of KB) can never match a pre-warmed prompt, and sending it as a
    PostgREST `?question=eq.<value>` URL filter overflows the request limit → a
    400 that surfaced as a 500 on the whole ask (the multi-file 'Failed to fetch'
    bug). find_cached_ask must short-circuit to a miss BEFORE touching the DB.

    The route-level tests can't reproduce this — the test DB enforces no URL
    limit — so we assert the guard directly: require_client must never be called.
    """
    from app.db import asks as asks_mod

    called = False

    def _boom():
        nonlocal called
        called = True
        raise AssertionError("require_client must not be called for an oversized question")

    original = asks_mod.require_client
    asks_mod.require_client = _boom
    try:
        huge = "q " * asks_mod._MAX_CACHE_QUESTION_CHARS  # well over the ceiling
        assert asks_mod.find_cached_ask("acme", huge) is None
        assert called is False
    finally:
        asks_mod.require_client = original


# ---- PRD-tab grounding (prd_id) ---------------------------------------------
# A chat running beside an open PRD sends prd_id; the answer must be grounded
# on that PRD (+ insight/evidence/tickets), the prd must be ownership-gated,
# and the (dataset, question)-keyed prewarm cache must be bypassed — it would
# serve a context-free answer for a question about the open PRD.

def _seed_prd(db, *, slug: str, prd_id: int, payload_md: str = "# PRD body"):
    brief = db.table("briefs").insert(
        {"dataset": slug, "week_label": "W",
         "payload": {"insights": [{"title": "Top insight", "body": "Insight body."}]},
         "is_current": True}
    ).execute().data[0]
    db.table("prds").insert(
        {"id": prd_id, "brief_id": brief["id"], "insight_index": 0,
         "title": "The open PRD", "status": "ready", "payload_md": payload_md}
    ).execute()


def test_ask_foreign_prd_id_returns_404(tenant_client, isolated_settings):
    """prd_id must belong to the caller — otherwise a crafted id would seed a
    foreign tenant's PRD into the answer context."""
    tenant_client.make(slug="company-a")
    _seed_prd(isolated_settings["supabase"], slug="company-a", prd_id=401)
    b = tenant_client.make(slug="company-b")
    _seed_corpus(isolated_settings["data_dir"], dataset="company-b")
    resp = b.client.post(
        "/v1/ask",
        json={"question": "What does this PRD say?", "dataset": "company-b",
              "prd_id": 401},
    )
    assert resp.status_code == 404


def test_ask_with_prd_id_grounds_answer_on_prd(
    tenant_client, isolated_settings, fake_llm
):
    """The LLM prompt for a PRD-tab ask carries the CURRENT PRD CONTEXT block
    with the PRD body and its source insight."""
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    _seed_prd(
        isolated_settings["supabase"], slug="acme", prd_id=402,
        payload_md="# Export revamp\nUsers need CSV export.",
    )
    fake_llm["payload"] = {
        "answer": "grounded", "key_points": [], "citations": [],
        "confidence": 0.9, "unanswered": "",
    }
    start = t.client.post(
        "/v1/ask",
        json={"question": "What is the biggest churn driver?",
              "dataset": "acme", "prd_id": 402},
    ).json()
    body = _poll_ask(t.client, start["ask_id"])
    assert body["status"] == "ready"
    assert len(fake_llm["calls"]) == 1
    prompt = fake_llm["calls"][0]["user"]
    assert "CURRENT PRD CONTEXT" in prompt
    assert "Users need CSV export." in prompt
    assert "Top insight" in prompt
    # The job row records the grounding PRD (mirrors conversation_id).
    row = db.get_ask_job(start["ask_id"])
    assert row["prd_id"] == 402


def test_ask_with_prd_id_skips_prewarm_cache(
    tenant_client, isolated_settings, fake_llm, monkeypatch
):
    """A cached (dataset, question) answer is context-free — a PRD-tab ask must
    generate fresh instead of serving it."""
    monkeypatch.setattr(ask_route, "CACHE_HIT_DELAY_MIN_SECONDS", 0.0)
    monkeypatch.setattr(ask_route, "CACHE_HIT_DELAY_MAX_SECONDS", 0.0)
    t = tenant_client.make(slug="acme")
    _seed_corpus(isolated_settings["data_dir"], dataset="acme")
    _seed_prd(isolated_settings["supabase"], slug="acme", prd_id=403)
    question = "What are the biggest revenue drivers"
    cache_id = db.start_cached_ask(dataset="acme", question=question)
    db.complete_cached_ask(
        cache_id,
        json.dumps({"answer": "**Cached answer**", "key_points": [],
                    "citations": [], "confidence": 1.0, "unanswered": ""}),
    )
    fake_llm["payload"] = {
        "answer": "fresh grounded answer", "key_points": [], "citations": [],
        "confidence": 0.9, "unanswered": "",
    }
    fake_llm["calls"].clear()
    start = t.client.post(
        "/v1/ask",
        json={"question": question, "dataset": "acme", "prd_id": 403},
    ).json()
    body = _poll_ask(t.client, start["ask_id"])
    assert body["status"] == "ready"
    assert body["answer"] == "fresh grounded answer"
    assert len(fake_llm["calls"]) == 1
