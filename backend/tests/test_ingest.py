"""Tests for app.ingest — file conversion + helpers.

docx/xlsx/pdf tests build real binary inputs with the upstream libs to make
sure our converters survive the round-trip. If a library isn't installed we
skip rather than fail — CI installs them via requirements.txt.
"""
from __future__ import annotations

import io

import pytest

from app import ingest


def test_slugify_basic():
    assert ingest.slugify("Acme Corp Quarterly Plan") == "acme_corp_quarterly_plan"
    assert ingest.slugify("UPPER") == "upper"
    assert ingest.slugify("with-hyphens") == "with_hyphens"
    assert ingest.slugify("   spaced  ") == "spaced"
    assert ingest.slugify("!!!") == "untitled"
    assert ingest.slugify("") == "untitled"


def test_md_filename_preserves_stem():
    assert ingest.md_filename("Customer Data.docx") == "customer_data.md"
    assert ingest.md_filename("path/to/Foo Bar.PDF") == "foo_bar.md"
    assert ingest.md_filename("already_a_slug.txt") == "already_a_slug.md"


def test_txt_to_md_passthrough():
    out = ingest.txt_to_md(b"hello\nworld")
    assert "hello" in out
    assert "world" in out


def test_txt_handles_non_utf8():
    out = ingest.txt_to_md(b"\xff\xfecaf\xe9")
    # Doesn't raise; replacement chars OK.
    assert isinstance(out, str)


def test_convert_unknown_textual_passes_through():
    # yaml/json/etc. have no dedicated converter but decode cleanly → text.
    out = ingest.convert("config.yaml", b"name: acme\nenv: prod\n")
    assert "name: acme" in out
    assert "env: prod" in out


def test_convert_unknown_binary_becomes_stub():
    # Binary content (e.g. audio) is stored but emits a placeholder stub.
    out = ingest.convert("memo.m4a", b"\x00\x01\x02binary\xff\xfe")
    assert "memo.m4a" in out
    assert "not yet parsed" in out


def test_convert_never_raises_for_unknown_type():
    # The old UnsupportedFileType path is gone — nothing is rejected.
    assert isinstance(ingest.convert("foo.exe", b"\x00\x01"), str)


def test_convert_routes_by_extension_case_insensitive():
    # .TXT and .txt both route to txt_to_md
    out = ingest.convert("notes.TXT", b"hi")
    assert out == "hi"


def test_docx_to_md():
    try:
        import docx
    except ImportError:
        pytest.skip("python-docx not installed")
    buf = io.BytesIO()
    d = docx.Document()
    d.add_heading("Title", level=1)
    d.add_paragraph("Body paragraph.")
    d.add_heading("Section", level=2)
    d.add_paragraph("Second body.")
    d.save(buf)
    out = ingest.docx_to_md(buf.getvalue())
    assert "# Title" in out
    assert "Body paragraph." in out
    assert "## Section" in out


def test_xlsx_to_md():
    try:
        import openpyxl
    except ImportError:
        pytest.skip("openpyxl not installed")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Funnel"
    ws.append(["step", "users"])
    ws.append(["signup", 100])
    ws.append(["activate", 60])
    buf = io.BytesIO()
    wb.save(buf)
    out = ingest.xlsx_to_md(buf.getvalue())
    assert "## Funnel" in out
    assert "signup" in out
    assert "100" in out


def test_xlsx_to_md_truncates_long_sheets():
    try:
        import openpyxl
    except ImportError:
        pytest.skip("openpyxl not installed")
    wb = openpyxl.Workbook()
    ws = wb.active
    for i in range(100):
        ws.append([f"row{i}", i])
    buf = io.BytesIO()
    wb.save(buf)
    out = ingest.xlsx_to_md(buf.getvalue(), max_rows=10)
    assert "truncated" in out


def test_pdf_to_md():
    try:
        import pypdf
        from pypdf import PdfWriter
    except ImportError:
        pytest.skip("pypdf not installed")
    # Easiest portable way: write an empty page (we don't need real text — we
    # only verify the converter doesn't crash and returns a string).
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    buf = io.BytesIO()
    writer.write(buf)
    out = ingest.pdf_to_md(buf.getvalue())
    assert isinstance(out, str)


def _build_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1: title + bullet body (Title and Content layout)
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Q3 Roadmap PRD"
    s1.placeholders[1].text = "Cut churn 20%\nShip SSO"
    # Slide 2: title only + a 2x2 table (Title Only layout)
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Metrics"
    tbl = s2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "Metric"
    tbl.cell(0, 1).text = "Target"
    tbl.cell(1, 0).text = "Churn"
    tbl.cell(1, 1).text = "8%"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_to_md_slides_and_tables():
    try:
        import pptx  # noqa: F401
    except ImportError:
        pytest.skip("python-pptx not installed")
    out = ingest.pptx_to_md(_build_pptx())
    # One section per slide, title + bullets preserved
    assert "## Slide 1" in out
    assert "Q3 Roadmap PRD" in out
    assert "Ship SSO" in out
    # Table rendered as pipe rows on slide 2
    assert "## Slide 2" in out
    assert "| Metric | Target |" in out
    assert "| Churn | 8% |" in out


def test_convert_routes_pptx():
    try:
        import pptx  # noqa: F401
    except ImportError:
        pytest.skip("python-pptx not installed")
    out = ingest.convert("deck.PPTX", _build_pptx())  # case-insensitive suffix
    assert "Q3 Roadmap PRD" in out
    assert ".pptx" in ingest.SUPPORTED_SUFFIXES


def test_convert_md_passthrough():
    out = ingest.convert("notes.md", b"# heading\n\nbody")
    assert out.startswith("# heading")


# ─────────────── Unreadable-file stubs are machine-detectable ───────────────
#
# The stub means "we could not read this file". It must never be mistaken for
# content: it used to reach the KG extractor and burn an LLM call mining
# signals out of its own apology text.


def test_stub_carries_the_marker_and_is_detected():
    out = ingest.convert("memo.m4a", b"\x00\x01\x02binary\xff\xfe")
    assert ingest.UNPARSED_STUB_MARKER in out
    assert ingest.is_unparsed_stub(out)


def test_real_content_is_not_mistaken_for_a_stub():
    for text in [
        ingest.convert("notes.txt", b"customers keep asking for SSO"),
        ingest.convert("config.yaml", b"name: acme\n"),
        "",
        None,
    ]:
        assert not ingest.is_unparsed_stub(text)


def test_legacy_stubs_without_the_marker_are_still_detected():
    """Files uploaded before the marker shipped must also be skipped."""
    legacy = (
        "# old.doc\n\n_Stored as a source but not yet parsed (type .doc, 12 KB). "
        "Binary or unrecognized format — its content is not included in "
        "analysis yet._\n"
    )
    assert ingest.is_unparsed_stub(legacy)


def test_marker_is_an_html_comment_so_it_never_renders():
    """It rides in the markdown, so it must be invisible when displayed."""
    assert ingest.UNPARSED_STUB_MARKER.startswith("<!--")
    assert ingest.UNPARSED_STUB_MARKER.endswith("-->")


# ─────────────────────────── HTML / RTF ───────────────────────────
#
# Both decode as text, so before this they "passed through" — the extractor
# received raw tags / control words as if they were customer evidence.


def test_html_is_read_as_text_not_raw_markup():
    html = (
        b"<html><head><style>.x{color:red}</style></head><body>"
        b"<nav>Home About</nav>"
        b"<h1>Customer feedback</h1>"
        b"<p>Users want SSO before renewal.</p>"
        b"<script>track()</script>"
        b"<footer>(c) 2026</footer>"
        b"</body></html>"
    )
    out = ingest.convert("export.html", html)
    assert "Customer feedback" in out
    assert "Users want SSO before renewal." in out
    # Markup and non-content chrome are gone.
    assert "<p>" not in out and "<h1>" not in out
    assert "track()" not in out
    assert "color:red" not in out


def test_htm_extension_routes_to_the_html_reader():
    out = ingest.convert("page.htm", b"<html><body><p>hello</p></body></html>")
    assert "hello" in out
    assert "<p>" not in out


def test_rtf_control_words_are_stripped():
    # Shaped like real RTF: \par is followed by a newline (a writer never runs
    # a control word straight into a word), and the font table is a NESTED
    # destination group.
    rtf = (
        rb"{\rtf1\ansi\deff0{\fonttbl{\f0 Times New Roman;}}" b"\n"
        rb"{\*\generator Riched20 10.0;}" b"\n"
        rb"\pard\f0\fs24 Churn risk on the Acme account.\par" b"\n"
        rb"Renewal is at risk.\par" b"\n"
        rb"}"
    )
    out = ingest.convert("note.rtf", rtf)
    assert "Churn risk on the Acme account." in out
    assert "Renewal is at risk." in out
    # None of the control layer or the metadata tables survive.
    assert "rtf1" not in out
    assert "fonttbl" not in out and "Times New Roman" not in out
    assert "generator" not in out and "Riched20" not in out
    assert r"\par" not in out


def test_rtf_keeps_the_word_after_a_paragraph_break():
    """Regression: a two-pass stripper that rewrote \\par before tokenizing
    control words swallowed the following word (`\\parRenewal` reads as one
    control word). Single-pass tokenizing is what makes this safe."""
    out = ingest.convert("note.rtf", rb"{\rtf1 First line.\par Second line.}")
    assert "First line." in out
    assert "Second line." in out


# ─────────────────── Tabular row caps are generous now ───────────────────
#
# 30 (xlsx) / 200 (csv) silently discarded most of a real analytics, revenue
# or CRM export — exactly the evidence-bearing categories.


def test_csv_keeps_far_more_than_the_old_200_row_cap():
    rows = [b"id,comment"] + [b"%d,needs SSO" % i for i in range(1, 1001)]
    out = ingest.convert("feedback.csv", b"\n".join(rows))
    assert ingest.CSV_MAX_ROWS >= 5000
    # A 1000-row export survives whole — no truncation notice.
    assert "truncated" not in out
    assert out.count("needs SSO") == 1000


def test_csv_still_truncates_a_pathological_file_and_says_so():
    n = ingest.CSV_MAX_ROWS + 10
    rows = [b"id,comment"] + [b"%d,x" % i for i in range(n)]
    out = ingest.convert("huge.csv", b"\n".join(rows))
    assert f"truncated to {ingest.CSV_MAX_ROWS}" in out


def test_tsv_passes_through_intact_rather_than_being_misparsed():
    """A comma-splitting CSV reader would collapse tab data into one column,
    so .tsv deliberately keeps the textual passthrough."""
    out = ingest.convert("export.tsv", b"id\tcomment\n1\tneeds SSO\n")
    assert "needs SSO" in out
